import os

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def create_poster():
    prs = Presentation()

    # Keep same poster dimensions.
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.26)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Palette
    NAVY = RGBColor(13, 27, 42)
    MID_BLUE = RGBColor(30, 70, 120)
    TEAL = RGBColor(42, 157, 143)
    BG = RGBColor(246, 249, 252)
    WHITE = RGBColor(255, 255, 255)
    TEXT = RGBColor(30, 36, 44)
    MUTED = RGBColor(78, 92, 108)
    BORDER = RGBColor(188, 200, 214)

    # Layout
    margin = 0.18
    gap = 0.10
    col_w = (11.69 - 2 * margin - 2 * gap) / 3.0
    x1 = margin
    x2 = x1 + col_w + gap
    x3 = x2 + col_w + gap
    y_start = 0.86
    y_gap = 0.08

    img_dir = os.path.join(os.path.dirname(__file__), "Ankit Chavda", "images", "img")

    def add_rect(left, top, width, height, fill, line=None, rounded=True):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = line or fill
        shape.line.width = Pt(0.8)
        return shape

    def add_text(left, top, width, height, lines, size=9, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(3)
        tf.margin_right = Pt(3)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)

        if isinstance(lines, str):
            lines = [lines]

        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = align
            p.space_after = Pt(0.5)
            for run in p.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = color
        return tb

    def add_panel(left, top, width, height, title, title_color):
        add_rect(left, top, width, height, WHITE, line=BORDER, rounded=True)
        add_rect(left, top, width, 0.18, title_color, line=title_color, rounded=False)
        add_text(left + 0.02, top + 0.01, width - 0.04, 0.16, title, size=8.6, bold=True, color=WHITE)
        inner_left = left + 0.03
        inner_top = top + 0.21
        inner_width = width - 0.06
        inner_height = height - 0.24
        return inner_left, inner_top, inner_width, inner_height

    def add_image_fit(path, left, top, width, height):
        if not os.path.exists(path):
            add_text(left, top, width, height, ["Image missing: " + os.path.basename(path)], size=8, color=MUTED)
            return
        try:
            with PILImage.open(path) as im:
                w, h = im.size
            if w <= 0 or h <= 0:
                add_text(left, top, width, height, ["Image invalid: " + os.path.basename(path)], size=8, color=MUTED)
                return
            ratio = min(width / w, height / h)
            out_w = w * ratio
            out_h = h * ratio
            x = left + (width - out_w) / 2.0
            y = top + (height - out_h) / 2.0
            slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(out_w), Inches(out_h))
        except Exception:
            add_text(left, top, width, height, ["Image load error: " + os.path.basename(path)], size=8, color=MUTED)

    # Background + header/footer
    add_rect(0, 0, 11.69, 8.26, BG, line=BG, rounded=False)
    add_rect(0, 0, 11.69, 0.72, NAVY, line=NAVY, rounded=False)
    add_rect(0, 0, 11.69, 0.03, TEAL, line=TEAL, rounded=False)
    add_rect(0, 7.90, 11.69, 0.36, NAVY, line=NAVY, rounded=False)
    add_rect(0, 7.90, 11.69, 0.03, TEAL, line=TEAL, rounded=False)

    add_text(
        0.20,
        0.06,
        7.60,
        0.25,
        "Hindi Verbal Fluency: VFT + SpAM Mental Lexicon Analysis",
        size=17,
        bold=True,
        color=WHITE,
    )
    add_text(
        0.20,
        0.34,
        8.80,
        0.22,
        "Team CTRL+ALT+DEL | Akshat Kotadia (2025201005) | Ankit Chavda (2025201045) | Om Mehra (2025201008)",
        size=9.6,
        color=RGBColor(214, 228, 247),
    )
    add_text(9.05, 0.10, 2.45, 0.42, "BRSM 2026 Poster", size=11, bold=True, color=RGBColor(220, 236, 255), align=PP_ALIGN.RIGHT)

    add_text(
        0.20,
        7.97,
        11.20,
        0.23,
        "VFT=Verbal Fluency Task | SpAM=Spatial Arrangement Method | IRT=Inter-Response Time | Semantic + Phonetic Similarity Both Included",
        size=8,
        color=RGBColor(220, 236, 255),
        align=PP_ALIGN.CENTER,
    )

    # ---------------- Column 1 ----------------
    c1_y1_h = 1.20
    c1_y2_h = 2.85
    c1_y3_h = 2.55

    p1 = add_panel(x1, y_start, col_w, c1_y1_h, "INTRODUCTION + DATASET", MID_BLUE)
    add_text(
        p1[0],
        p1[1],
        p1[2],
        p1[3],
        [
            "- 35 Hindi-English bilingual participants",
            "- VFT: 60-sec word generation in animals/foods/colours/body-parts",
            "- SpAM: participants arranged words by similarity in 2D",
            "- Analysis subset: 723 Hindi/Hinglish responses",
            "- Colours coverage low (n=4), used mainly for descriptive reading",
        ],
        size=8.9,
    )

    c1_y2 = y_start + c1_y1_h + y_gap
    p2 = add_panel(x1, c1_y2, col_w, c1_y2_h, "RESEARCH QUESTIONS + HYPOTHESES", TEAL)
    add_text(
        p2[0],
        p2[1],
        p2[2],
        1.14,
        [
            "RQ1 domain effect on words/IRT; RQ2 participant-level effects.",
            "RQ3 SpAM compactness; RQ4 semantic vs phonetic alignment.",
            "RQ5 integrated score vs VFT-only; RQ6-RQ7 phonology effects.",
            "",
            "H1 semantic clustering; H2 lexical exhaustion; H3-H4 domain differences.",
            "H5-H7 confidence links; H8 compactness; H9 semantic ARI; H10 phonetic ARI.",
            "H11 serial similarity shift; H12 phonology-productivity; H13 score comparison.",
            "",
            "All hypotheses H1-H13 are included in detailed results table (column 2).",
        ],
        size=8.1,
    )

    c1_y3 = c1_y2 + c1_y2_h + y_gap
    p3 = add_panel(x1, c1_y3, col_w, c1_y3_h, "METHODOLOGY + TRANSFORMERS", MID_BLUE)
    add_text(
        p3[0],
        p3[1],
        p3[2],
        p3[3],
        [
            "1) Merge VFT + SpAM + survey; keep Hindi/Hinglish rows.",
            "2) Compute total words, IRT, cluster size, and switches.",
            "3) Normality checks (Shapiro), then suitable inferential tests.",
            "4) H1 Welch t; H2/H11 mixed models; H3/H4/H8 Kruskal.",
            "5) H5-H7/H12 Spearman; H9/H10 ARI + sign/permutation; H13 Steiger.",
            "",
            "Transformer/model names and purpose:",
            "- paraphrase-multilingual-MiniLM-L12-v2: semantic embedding space (384-d)",
            "- phonetic-key representation: sound similarity signal",
            "- KMeans/Agglomerative/HDBSCAN + ARI/NMI for structure and alignment",
        ],
        size=8.4,
    )

    # ---------------- Column 2 ----------------
    c2_y1_h = 3.40
    c2_y2_h = 0.65
    c2_y3_h = 1.20
    c2_y4_h = 1.20

    p4 = add_panel(x2, y_start, col_w, c2_y1_h, "RESULTS TABLE (ALL H1-H13)", MID_BLUE)

    add_text(p4[0], p4[1], p4[2], 0.20, "Notebook-derived values only; no fabricated numbers.", size=7.9, color=MUTED)

    table_top = p4[1] + 0.20
    table_h = c2_y1_h - 0.48
    table = slide.shapes.add_table(14, 5, Inches(p4[0]), Inches(table_top), Inches(p4[2]), Inches(table_h)).table

    col_widths = [0.28, 0.27, 0.62, 1.83, 0.57]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    headers = ["RQ", "H", "Test", "Main Result", "Decision"]
    data_rows = [
        ["Found.", "H1", "Welch t", "t=9.30, p<0.001, d=1.12", "Support"],
        ["Found.", "H2", "MixedLM", "No positive serial slope (0/4)", "No"],
        ["RQ1", "H3", "Kruskal", "H=3.2639, p=0.1956", "No"],
        ["RQ1", "H4", "Kruskal", "H=3.9028, p=0.1421", "No"],
        ["RQ2", "H5", "Spearman", "rho=-0.3951, p(one)=0.9906", "No"],
        ["RQ2", "H6", "Spearman", "rho=+0.2755, p(one)=0.9454", "No"],
        ["RQ2", "H7", "Spearman", "rho=-0.2259, p(one)=0.8969", "No"],
        ["RQ3", "H8", "Kruskal", "H=3.4451, p=0.3280", "No"],
        ["RQ4", "H9", "ARI+Sign", "mean ARI=0.1765, sign p<0.0001", "Mixed"],
        ["RQ4", "H10", "ARI+Sign", "mean ARI=0.0745, p=0.0662", "No"],
        ["RQ6", "H11", "LME", "beta=-0.004144, p=0.1028", "No"],
        ["RQ7", "H12", "Spearman", "rho=+0.114, p=0.2570", "No"],
        ["RQ5", "H13", "Steiger", "z=0.724, p=0.4689", "No diff"],
    ]

    row_h = table_h / 14.0
    for r in range(14):
        table.rows[r].height = Inches(row_h)

    def style_cell(cell, text, size=7.2, bold=False, color=TEXT, align=PP_ALIGN.LEFT, fill=None):
        if fill is not None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color

    for c in range(5):
        style_cell(table.cell(0, c), headers[c], size=7.4, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=MID_BLUE)

    for r, row in enumerate(data_rows, start=1):
        fill = RGBColor(246, 249, 252) if r % 2 == 0 else WHITE
        for c in range(5):
            style_cell(
                table.cell(r, c),
                row[c],
                size=7.1,
                bold=False,
                color=TEXT,
                align=PP_ALIGN.CENTER if c in (0, 1, 4) else PP_ALIGN.LEFT,
                fill=fill,
            )

    c2_y2 = y_start + c2_y1_h + y_gap
    p5 = add_panel(x2, c2_y2, col_w, c2_y2_h, "KEY FINDINGS", TEAL)
    add_text(
        p5[0],
        p5[1],
        p5[2],
        p5[3],
        [
            "H1 is strong and clear; most other hypotheses are not supported in this sample.",
            "Semantic signal is stronger than phonetic signal in current data.",
        ],
        size=8.0,
    )

    c2_y3 = c2_y2 + c2_y2_h + y_gap
    p6 = add_panel(x2, c2_y3, col_w, c2_y3_h, "IMPORTANT FIGURE 1 (H1)", MID_BLUE)
    img1 = os.path.join(img_dir, "fig_moduleA_h1_cluster_comparison.png")
    add_image_fit(img1, p6[0], p6[1], p6[2], 0.75)
    add_text(p6[0], p6[1] + 0.76, p6[2], 0.18, "Interpretation: between-cluster pause is clearly higher than within-cluster pause.", size=7.7, color=MUTED)

    c2_y4 = c2_y3 + c2_y3_h + y_gap
    p7 = add_panel(x2, c2_y4, col_w, c2_y4_h, "IMPORTANT FIGURE 2 (H9/H10)", TEAL)
    img2 = os.path.join(img_dir, "fig_moduleD_h9_h10_alignment.png")
    add_image_fit(img2, p7[0], p7[1], p7[2], 0.75)
    add_text(p7[0], p7[1] + 0.76, p7[2], 0.18, "Interpretation: semantic ARI tends to be above phonetic ARI, but effect is modest.", size=7.7, color=MUTED)

    # ---------------- Column 3 ----------------
    c3_y1_h = 1.05
    c3_y2_h = 1.45
    c3_y3_h = 1.45
    c3_y4_h = 2.50

    p8 = add_panel(x3, y_start, col_w, c3_y1_h, "SEMANTIC VS PHONETIC SUMMARY", MID_BLUE)
    add_text(
        p8[0],
        p8[1],
        p8[2],
        p8[3],
        [
            "Semantic similarity: weak-to-moderate structure present (H9 mixed).",
            "Phonetic similarity: no clear confirmatory support (H10/H12 not supported).",
        ],
        size=8.3,
    )

    c3_y2 = y_start + c3_y1_h + y_gap
    p9 = add_panel(x3, c3_y2, col_w, c3_y2_h, "IMPORTANT FIGURE 3 (H11/H12)", TEAL)
    img3 = os.path.join(img_dir, "graph_12_1_h11_dual_similarity.png")
    add_image_fit(img3, p9[0], p9[1], p9[2], 0.92)
    add_text(p9[0], p9[1] + 0.95, p9[2], 0.18, "Interpretation: no reliable phonological facilitation trend over retrieval order.", size=7.7, color=MUTED)

    c3_y3 = c3_y2 + c3_y2_h + y_gap
    p10 = add_panel(x3, c3_y3, col_w, c3_y3_h, "IMPORTANT FIGURE 4 (H13)", MID_BLUE)
    img4 = os.path.join(img_dir, "graph_13_1_composite_vs_vft_confidence.png")
    add_image_fit(img4, p10[0], p10[1], p10[2], 0.92)
    add_text(p10[0], p10[1] + 0.95, p10[2], 0.18, "Interpretation: integrated score does not show significant gain over VFT-only score.", size=7.7, color=MUTED)

    c3_y4 = c3_y3 + c3_y3_h + y_gap
    p11 = add_panel(x3, c3_y4, col_w, c3_y4_h, "DISCUSSION + CONCLUSION", TEAL)
    add_text(
        p11[0],
        p11[1],
        p11[2],
        p11[3],
        [
            "- Strongest and most reliable result is semantic clustering boundary cost (H1).",
            "- Domain differences and confidence-linked effects were not statistically reliable.",
            "- Sample size and confidence ceiling can suppress detectable associations.",
            "- Semantic organisation appears clearer than phonetic organisation here.",
            "",
            "Final conclusion: Hindi lexical retrieval is semantically structured,",
            "but most secondary hypotheses remain unsupported in this dataset.",
        ],
        size=8.7,
    )

    output = os.path.join(os.path.dirname(__file__), "Hindi_Fluency_Poster_Clean.pptx")
    prs.save(output)
    print("DONE:", output)


if __name__ == "__main__":
    create_poster()
